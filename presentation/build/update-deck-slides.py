#!/usr/bin/env python3
"""Update Voyage Canvas demo deck: title slide, architecture fixes, rollout + thank-you."""

from __future__ import annotations

import copy
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
DECK_PATHS = [
    ROOT / "presentation" / "Voyage_Canvas_RCG_Demo.pptx",
    ROOT / "Voyage_Canvas_RCG_Demo.pptx",
]

SLIDE1_EYEBROW = "Royal Caribbean"
SLIDE1_TITLE = "Voyage Canvas"
SLIDE1_BYLINE = "Manoj Pammina · August 17, 2026"
SLIDE1_PROMPT = "Opening: end-to-end AI cruise planning demo."
SLIDE1_PROVES = (
    "A guest-facing assistant can be AI-forward while deterministic services "
    "keep commerce truth authoritative."
)

TOOLS_TEXT = (
    "search_sailings\n"
    "check_availability\n"
    "get_pricing\n"
    "get_policy_content\n"
    "create_hold\n"
    "start_booking"
)

SLIDE17_EYEBROW = "Architecture route + product proof"
SLIDE17_TITLE = "Close: launch the assistant safely"
SLIDE17_PROMPT = "Closing: staged rollout and the next hardening steps before broad guest exposure."
SLIDE17_PROVES = (
    "The solution is a working vertical slice with a clear production path "
    "and explicit safety boundaries."
)
SLIDE17_BULLETS = [
    "Demo path: natural language -> verified sailings -> policy RAG -> hold -> signed checkout handoff",
    "Safety path: model failure falls back to deterministic guided planning",
    "Authority boundary: Gemini explains; services decide price, availability, holds, and booking handoff",
    "Launch path: employee alpha -> shadow mode -> feature-flagged beta -> monitored rollout",
    "Next hardening: real API adapters, managed vector store, broader validators, scheduled reconciliation",
]

SLIDE18_EYEBROW = "Royal Caribbean · Voyage Canvas"
SLIDE18_TITLE = "Thank you"
SLIDE18_MESSAGE = "Questions welcome."
SLIDE18_FOOTER = "Manoj Pammina · August 17, 2026"


def set_shape_text(shape, text: str) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.text = text


def style_paragraph(shape, *, size: int = 18, bold: bool = False) -> None:
    for paragraph in shape.text_frame.paragraphs:
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold


def remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def find_text_shapes(slide):
    return [shape for shape in slide.shapes if shape.has_text_frame and shape.text_frame.text.strip()]


def add_centered_textbox(
    slide,
    *,
    slide_width: int,
    text: str,
    top: int,
    width: int,
    height: int,
    size: int,
    bold: bool = False,
) -> None:
    left = (slide_width - width) // 2
    box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    text_frame.word_wrap = True


def duplicate_slide(presentation, index: int):
    source = presentation.slides[index]
    destination = presentation.slides.add_slide(source.slide_layout)
    for shape in source.shapes:
        new_element = copy.deepcopy(shape.element)
        destination.shapes._spTree.insert_element_before(new_element, "p:extLst")
    return destination


def move_slide(presentation, old_index: int, new_index: int) -> None:
    slides = presentation.slides._sldIdLst
    entries = list(slides)
    entry = entries[old_index]
    slides.remove(entry)
    slides.insert(new_index, entry)


def update_footers(presentation, total_slides: int) -> None:
    pattern = re.compile(r"Voyage Canvas demo route \| \d+/\d+")
    replacement_template = "Voyage Canvas demo route | {}/{}"
    for index, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if pattern.fullmatch(text):
                set_shape_text(shape, replacement_template.format(index + 1, total_slides))


def update_slide_one(slide) -> None:
    text_shapes = find_text_shapes(slide)
    if len(text_shapes) >= 2:
        set_shape_text(text_shapes[0], SLIDE1_EYEBROW)
        set_shape_text(text_shapes[1], SLIDE1_TITLE)

    byline = None
    for shape in text_shapes:
        text = shape.text_frame.text.strip()
        if "Manoj Pammina" in text or "August 17" in text:
            byline = shape
            break

    if byline is None:
        byline = slide.shapes.add_textbox(
            Inches(8.8),
            Inches(0.31),
            Inches(3.8),
            Inches(0.44),
        )

    set_shape_text(byline, SLIDE1_BYLINE)
    style_paragraph(byline, size=16, bold=True)

    labels = [shape.text_frame.text.strip() for shape in text_shapes]
    if "Demo prompt / action" in labels:
        idx = labels.index("Demo prompt / action")
        if idx + 1 < len(text_shapes):
            set_shape_text(text_shapes[idx + 1], SLIDE1_PROMPT)
    if "What this proves" in labels:
        idx = labels.index("What this proves")
        if idx + 1 < len(text_shapes):
            set_shape_text(text_shapes[idx + 1], SLIDE1_PROVES)


def fix_architecture_layout(slide) -> None:
    tools_box = None
    checkout_box = None
    tools_title = None
    tools_body = None
    checkout_title = None
    checkout_body = None
    handoff = None

    for shape in slide.shapes:
        if shape.shape_type == 1 and shape.left == 9010650:
            if 3700000 < shape.top < 4500000:
                tools_box = shape
            elif shape.top > 5000000:
                checkout_box = shape
            continue
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text == "Deterministic tools":
            tools_title = shape
        elif text.startswith("search_sailings"):
            tools_body = shape
        elif text == "Existing checkout":
            checkout_title = shape
        elif text.startswith("Signed booking context"):
            checkout_body = shape
        elif "handoff" in text:
            handoff = shape

    if tools_box is not None:
        tools_box.top = 3810000
        tools_box.height = 1660000
    if tools_title is not None:
        tools_title.top = 3890000
    if tools_body is not None:
        tools_body.top = 4340000
        tools_body.height = 820000
        set_shape_text(tools_body, TOOLS_TEXT)
        for paragraph in tools_body.text_frame.paragraphs:
            paragraph.font.size = Pt(9)

    checkout_top = 5650650
    if checkout_box is not None:
        checkout_box.top = checkout_top
        checkout_box.height = 960000
    if checkout_title is not None:
        checkout_title.top = checkout_top + 40000
    if checkout_body is not None:
        checkout_body.top = checkout_top + 440000
        checkout_body.height = 420000

    # Horizontal connector captions need the full gap between their two boxes,
    # otherwise the arrow glyph wraps onto a second line and spills out.
    connector_bounds = {
        "SSE actions →": (2207950, 1132950),
        "typed request →": (4853000, 1326350),
    }
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        bounds = connector_bounds.get(shape.text_frame.text.strip())
        if bounds is None:
            continue
        shape.left, shape.width = bounds
        shape.text_frame.word_wrap = True
        style_paragraph(shape, size=11, bold=True)

    if handoff is not None:
        set_shape_text(handoff, "signed\nhandoff ↓")
        if checkout_box is not None:
            handoff.left = checkout_box.left - 520000
            handoff.top = checkout_box.top - 320000
        handoff.width = 720000
        handoff.height = 420000
        style_paragraph(handoff, size=12, bold=True)


def fix_rollout_slide_layout(slide) -> None:
    eyebrow = None
    title = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text == SLIDE17_EYEBROW:
            eyebrow = shape
        elif text == SLIDE17_TITLE:
            title = shape

    if eyebrow is not None:
        eyebrow.top = 190500
        eyebrow.height = 228600
        eyebrow.width = 3048000
        style_paragraph(eyebrow, size=12)
    if title is not None:
        title.top = 400050
        title.height = 419100
        title.width = 10668000
        style_paragraph(title, size=26.25, bold=True)

    bullet_shapes = [
        shape
        for shape in slide.shapes
        if shape.has_text_frame and shape.text_frame.text.strip().startswith("•")
    ]
    bullet_top_start = 1676400
    bullet_gap = 660400
    bullet_height = 560000

    for index, bullet in enumerate(SLIDE17_BULLETS):
        if index < len(bullet_shapes):
            shape = bullet_shapes[index]
        else:
            shape = slide.shapes.add_textbox(4171950, bullet_top_start + index * bullet_gap, 7086600, bullet_height)
            bullet_shapes.append(shape)
        shape.top = bullet_top_start + index * bullet_gap
        shape.height = bullet_height
        shape.width = 7086600
        shape.left = 4171950
        set_shape_text(shape, f"• {bullet}")
        shape.text_frame.word_wrap = True
        style_paragraph(shape, size=15)


def update_architecture_slide(slide) -> None:
    fix_architecture_layout(slide)


def set_bullet_slide_content(slide, *, eyebrow, title, prompt, proves, bullets) -> None:
    text_shapes = find_text_shapes(slide)
    if len(text_shapes) >= 2:
        set_shape_text(text_shapes[0], eyebrow)
        set_shape_text(text_shapes[1], title)

    labels = [shape.text_frame.text.strip() for shape in text_shapes]
    if "Demo prompt / action" in labels:
        idx = labels.index("Demo prompt / action")
        if idx + 1 < len(text_shapes):
            set_shape_text(text_shapes[idx + 1], prompt)
    if "What this proves" in labels:
        idx = labels.index("What this proves")
        if idx + 1 < len(text_shapes):
            set_shape_text(text_shapes[idx + 1], proves)

    bullet_shapes = [shape for shape in text_shapes if shape.text_frame.text.strip().startswith("•")]
    for bullet_shape, bullet in zip(bullet_shapes, bullets):
        set_shape_text(bullet_shape, f"• {bullet}")
    for bullet_shape in bullet_shapes[len(bullets) :]:
        set_shape_text(bullet_shape, "")


def rebuild_thank_you_slide(slide, slide_width: int, slide_height: int) -> None:
    for shape in list(slide.shapes)[2:]:
        remove_shape(shape)

    add_centered_textbox(
        slide,
        slide_width=slide_width,
        text=SLIDE18_EYEBROW,
        top=int(slide_height * 0.18),
        width=int(slide_width * 0.82),
        height=int(slide_height * 0.08),
        size=22,
        bold=True,
    )
    add_centered_textbox(
        slide,
        slide_width=slide_width,
        text=SLIDE18_TITLE,
        top=int(slide_height * 0.36),
        width=int(slide_width * 0.82),
        height=int(slide_height * 0.18),
        size=56,
        bold=True,
    )
    add_centered_textbox(
        slide,
        slide_width=slide_width,
        text=SLIDE18_MESSAGE,
        top=int(slide_height * 0.56),
        width=int(slide_width * 0.72),
        height=int(slide_height * 0.1),
        size=30,
    )
    add_centered_textbox(
        slide,
        slide_width=slide_width,
        text=SLIDE18_FOOTER,
        top=int(slide_height * 0.86),
        width=int(slide_width * 0.82),
        height=int(slide_height * 0.07),
        size=16,
    )

    notes = slide.notes_slide
    if notes is not None:
        notes.notes_text_frame.text = (
            "Thank the panel and invite questions.\n\n"
            "If asked for the repo or runbook, mention github.com/manojpammina1/voyage-canvas "
            "after the session."
        )


def has_rollout_slide(presentation) -> bool:
    if len(presentation.slides) < 17:
        return False
    for shape in presentation.slides[16].shapes:
        if shape.has_text_frame and "launch the assistant safely" in shape.text_frame.text:
            return True
    return False


def update_deck(path: Path) -> None:
    presentation = Presentation(str(path))
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    update_slide_one(presentation.slides[0])
    update_architecture_slide(presentation.slides[1])

    if not has_rollout_slide(presentation):
        thank_you_index = len(presentation.slides) - 1
        duplicate_slide(presentation, thank_you_index - 1)
        rollout_index = len(presentation.slides) - 1
        move_slide(presentation, rollout_index, thank_you_index)

        set_bullet_slide_content(
            presentation.slides[thank_you_index],
            eyebrow=SLIDE17_EYEBROW,
            title=SLIDE17_TITLE,
            prompt=SLIDE17_PROMPT,
            proves=SLIDE17_PROVES,
            bullets=SLIDE17_BULLETS,
        )

        rollout_notes = presentation.slides[thank_you_index].notes_slide
        if rollout_notes is not None:
            rollout_notes.notes_text_frame.text = (
                "Cover P1-ROLL explicitly: employee alpha, shadow mode, feature-flagged beta, "
                "monitored rollout. Name scheduled reconciliation as a next hardening step."
            )

    if len(presentation.slides) >= 17:
        fix_rollout_slide_layout(presentation.slides[16])

    rebuild_thank_you_slide(
        presentation.slides[-1],
        slide_width,
        slide_height,
    )

    update_footers(presentation, len(presentation.slides))
    presentation.save(str(path))
    print(f"Updated {path} ({len(presentation.slides)} slides)")


def main() -> int:
    targets = [Path(arg) for arg in sys.argv[1:]] or DECK_PATHS
    for path in targets:
        if not path.exists():
            print(f"Skip missing deck: {path}", file=sys.stderr)
            continue
        update_deck(path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
